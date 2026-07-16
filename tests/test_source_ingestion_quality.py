from __future__ import annotations

import json
import os
import shutil
import subprocess
import sys
import tempfile
import unittest
from pathlib import Path
from unittest import mock

import yaml
from docx import Document
from pptx import Presentation
from pypdf import PdfWriter

from tests.test_support import load_script_module


REPO_ROOT = Path(__file__).resolve().parents[1]
SCRIPTS = REPO_ROOT / "00_system" / "scripts"
FIXTURE = REPO_ROOT / "tests" / "fixtures" / "ingestion" / "long-guide.md"
quality_module = load_script_module(SCRIPTS / "source_quality.py", "source_quality_test_module")
ingest_module = load_script_module(SCRIPTS / "ingest_source.py", "ingest_source_test_module")


def frontmatter(path: Path) -> dict[str, object]:
    text = path.read_text(encoding="utf-8")
    _, raw, _ = text.split("---", 2)
    payload = yaml.safe_load(raw)
    return payload if isinstance(payload, dict) else {}


class SourceQualityUnitTests(unittest.TestCase):
    def test_missing_source_is_not_misreported_as_an_ocr_problem(self) -> None:
        with tempfile.TemporaryDirectory() as temp:
            with self.assertRaises(FileNotFoundError):
                quality_module.audit_source(Path(temp) / "missing.pdf")

    def test_blank_pdf_is_blocked_and_requires_ocr(self) -> None:
        with tempfile.TemporaryDirectory() as temp:
            path = Path(temp) / "scan.pdf"
            writer = PdfWriter()
            for _ in range(4):
                writer.add_blank_page(width=612, height=792)
            with path.open("wb") as handle:
                writer.write(handle)
            quality = quality_module.audit_source(path)
            self.assertEqual(quality.status, "blocked")
            self.assertEqual(quality.total_units, 4)
            self.assertEqual(quality.extracted_units, 0)
            self.assertTrue(quality.needs_ocr)

    def test_partially_empty_presentation_requires_review(self) -> None:
        with tempfile.TemporaryDirectory() as temp:
            path = Path(temp) / "slides.pptx"
            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[1])
            slide.shapes.title.text = "Architecture"
            slide.placeholders[1].text = "A" * 240
            presentation.slides.add_slide(presentation.slide_layouts[6])
            presentation.save(path)
            quality = quality_module.audit_source(path)
            self.assertEqual(quality.status, "review")
            self.assertEqual(quality.unit_coverage, 0.5)

    def test_docx_uses_python_fallback_when_pandoc_is_unavailable(self) -> None:
        with tempfile.TemporaryDirectory() as temp:
            path = Path(temp) / "guide.docx"
            document = Document()
            document.add_heading("Architecture", level=1)
            document.add_paragraph("The source map is created before section notes and every section keeps provenance.")
            document.save(path)
            with mock.patch.object(ingest_module.subprocess, "run", side_effect=FileNotFoundError):
                extracted = ingest_module.extract_text_from_docx(path)
            self.assertIn("# Architecture", extracted)
            self.assertIn("keeps provenance", extracted)

    def test_pdf_blocks_follow_detected_chapter_boundaries(self) -> None:
        pages = []
        for number in range(1, 13):
            heading = "第1章 Architecture" if number == 1 else "第2章 Workflow" if number == 7 else ""
            pages.append(f"## 第 {number} 页\n\n{heading}\n" + (f"page {number} content " * 40))
        sections = ingest_module.build_source_sections("\n\n".join(pages), "pdf")
        self.assertEqual(len(sections), 2)
        self.assertEqual(sections[0].refs, [f"p.{number}" for number in range(1, 7)])
        self.assertEqual(sections[1].refs, [f"p.{number}" for number in range(7, 13)])

    def test_single_digit_procedure_steps_do_not_become_pdf_chapters(self) -> None:
        pages = [
            "## 第 1 页\n\n01 Architecture\n" + ("architecture content " * 30),
            "## 第 2 页\n\n1 Install dependency\n" + ("installation content " * 30),
            "## 第 3 页\n\n2 Start service\n" + ("service content " * 30),
        ]
        sections = ingest_module.build_source_sections("\n\n".join(pages), "pdf")
        self.assertEqual(len(sections), 1)
        self.assertEqual(sections[0].refs, ["p.1", "p.2", "p.3"])

    def test_pdf_front_matter_and_toc_are_grouped_before_body_sections(self) -> None:
        pages = [
            "## 第 1 页\n\nDocument title\nVersion 1.0",
            "## 第 2 页\n\n目录\nTable of Contents\n01 Architecture\n02 Workflow",
            "## 第 3 页\n\n03 Operations\n04 Security\nAppendix A FAQ",
            "## 第 4 页\n\n01 Architecture\nArchitecture\n这是正文内容。" + ("架构说明" * 30),
        ]
        sections = ingest_module.build_source_sections("\n\n".join(pages), "pdf")
        self.assertEqual(sections[0].title, "封面与目录")
        self.assertEqual(sections[0].refs, ["p.1", "p.2", "p.3"])
        self.assertEqual(sections[1].title, "01 Architecture")
        self.assertEqual(sections[1].refs, ["p.4"])

    def test_pdf_numeric_facts_and_code_comments_do_not_become_chapters(self) -> None:
        pages = [
            "## 第 1 页\n\n01 Architecture\nArchitecture\n" + ("architecture content " * 30),
            "## 第 2 页\n\nMetrics\n43 lines of code are enough for this example.\n# Update stable version\n" + ("metrics content " * 30),
            "## 第 3 页\n\n02 Workflow\nWorkflow\n" + ("workflow content " * 30),
        ]
        sections = ingest_module.build_source_sections("\n\n".join(pages), "pdf")
        self.assertEqual([section.title for section in sections], ["01 Architecture", "02 Workflow"])
        self.assertEqual(sections[0].refs, ["p.1", "p.2"])
        self.assertEqual(sections[1].refs, ["p.3"])

    def test_long_pdf_chapters_keep_heading_on_continuation_sections(self) -> None:
        pages = []
        for number in range(1, 15):
            heading = "01 Architecture\nArchitecture" if number == 1 else "02 Workflow\nWorkflow" if number == 8 else ""
            pages.append(f"## 第 {number} 页\n\n{heading}\n" + (f"page {number} content " * 40))
        sections = ingest_module.build_source_sections("\n\n".join(pages), "pdf")
        self.assertEqual(
            [section.title for section in sections],
            ["01 Architecture", "01 Architecture（续 2）", "02 Workflow", "02 Workflow（续 2）"],
        )
        self.assertEqual(sections[1].refs, ["p.7"])
        self.assertEqual(sections[3].refs, ["p.14"])

    def test_english_sentence_is_not_mistaken_for_a_lettered_appendix(self) -> None:
        pages = [
            "## 第 1 页\n\n01 Architecture\nArchitecture\n" + ("architecture content " * 30),
            "## 第 2 页\n\nA system design principle\n" + ("design content " * 30),
            "## 第 3 页\n\nA 常见问题 FAQ\nFrequently Asked Questions\n" + ("answer content " * 30),
        ]
        sections = ingest_module.build_source_sections("\n\n".join(pages), "pdf")
        self.assertEqual([section.title for section in sections], ["01 Architecture", "A 常见问题 FAQ"])
        self.assertEqual(sections[0].refs, ["p.1", "p.2"])
        self.assertEqual(sections[1].refs, ["p.3"])

    def test_pdf_cjk_spacing_is_normalized_and_excerpt_limit_is_exact(self) -> None:
        self.assertEqual(quality_module.normalize_pdf_text("国 际 平 台 OpenClaw"), "国际平台 OpenClaw")
        excerpt = ingest_module.excerpt_from_text("正文" * 1000, limit=1200)
        self.assertLessEqual(len(excerpt), 1200)
        self.assertIn("摘录已截断", excerpt)


class SourceIngestionGoldenTests(unittest.TestCase):
    def setUp(self) -> None:
        self.temp_dir = tempfile.TemporaryDirectory()
        self.root = Path(self.temp_dir.name)
        self.vault = self.root / "vault"
        self.vault.mkdir()
        shutil.copytree(REPO_ROOT / "00_system" / "templates", self.vault / "00_system" / "templates")
        registry = self.vault / "00_system" / "registry"
        registry.mkdir(parents=True)
        for name in ("page_schemas.json", "ingestion_quality.json"):
            shutil.copy2(REPO_ROOT / "00_system" / "registry" / name, registry / name)
        self.env = os.environ.copy()
        self.env["OBSIDIAN_WIKI_ROOT"] = str(self.vault)
        self.env["PYTHONIOENCODING"] = "utf-8"

    def tearDown(self) -> None:
        self.temp_dir.cleanup()

    def ingest(self, source: Path, title: str) -> Path:
        result = subprocess.run(
            [sys.executable, str(SCRIPTS / "ingest_source.py"), "--source", str(source), "--title", title],
            cwd=REPO_ROOT,
            env=self.env,
            check=False,
            capture_output=True,
            text=True,
            encoding="utf-8",
        )
        self.assertEqual(result.returncode, 0, result.stderr)
        return Path(result.stdout.strip().splitlines()[-1])

    def test_golden_long_document_creates_map_then_section_index_with_quality(self) -> None:
        note = self.ingest(FIXTURE, "Golden Long Guide")
        metadata = frontmatter(note)
        self.assertEqual(metadata["quality_status"], "pass")
        derived = metadata["derived_pages"]
        self.assertTrue(derived[0].endswith("-document-map.md"))
        self.assertGreaterEqual(len(derived), 4)
        document_map = self.vault / derived[0]
        map_metadata = frontmatter(document_map)
        self.assertEqual(map_metadata["section_count"], len(map_metadata["derived_sections"]))
        self.assertEqual(map_metadata["quality_status"], "pass")
        first_section = self.vault / map_metadata["derived_sections"][0]
        self.assertTrue(frontmatter(first_section)["source_refs"])

    def test_blocked_pdf_preserves_source_but_creates_no_weak_derivatives(self) -> None:
        source = self.root / "blank.pdf"
        writer = PdfWriter()
        writer.add_blank_page(width=612, height=792)
        writer.add_blank_page(width=612, height=792)
        with source.open("wb") as handle:
            writer.write(handle)
        note = self.ingest(source, "Blank Scan")
        metadata = frontmatter(note)
        self.assertEqual(metadata["quality_status"], "blocked")
        self.assertTrue(metadata["needs_ocr"])
        self.assertEqual(metadata["derived_pages"], [])
        self.assertTrue((self.vault / metadata["source_path"]).exists())

    def test_reingestion_removes_only_obsolete_generated_sections(self) -> None:
        source = self.root / "changing-guide.md"
        source.write_text("# One\n\nFirst body.\n\n# Two\n\nSecond body.\n\n# Three\n\nThird body.\n", encoding="utf-8")
        first_note = self.ingest(source, "Changing Guide")
        first_map = self.vault / frontmatter(first_note)["derived_pages"][0]
        first_sections = frontmatter(first_map)["derived_sections"]
        self.assertEqual(len(first_sections), 3)

        source.write_text("# One\n\nFirst body.\n\n# Two\n\nSecond body.\n", encoding="utf-8")
        second_note = self.ingest(source, "Changing Guide")
        second_map = self.vault / frontmatter(second_note)["derived_pages"][0]
        second_sections = frontmatter(second_map)["derived_sections"]
        self.assertEqual(len(second_sections), 2)
        self.assertFalse((self.vault / first_sections[-1]).exists())
        self.assertEqual(len(list((second_map.parent / "changing-guide-sections").glob("*.md"))), 2)


if __name__ == "__main__":
    unittest.main()
