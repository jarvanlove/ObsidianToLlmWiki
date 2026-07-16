from __future__ import annotations

import argparse
import json
import re
import unicodedata
from dataclasses import asdict, dataclass
from pathlib import Path

from docx import Document
from pptx import Presentation
from pypdf import PdfReader


SCRIPT_DIR = Path(__file__).resolve().parent
QUALITY_CONFIG_PATH = SCRIPT_DIR.parent / "registry" / "ingestion_quality.json"
TEXT_EXTENSIONS = {".md", ".markdown", ".txt", ".py", ".js", ".ts", ".json", ".yaml", ".yml", ".html", ".css", ".java", ".go", ".rs", ".sql"}


@dataclass(frozen=True)
class SourceQuality:
    quality_version: int
    extract_mode: str
    status: str
    total_units: int
    extracted_units: int
    empty_units: int
    extracted_chars: int
    unit_coverage: float
    average_chars_per_extracted_unit: float
    garbled_ratio: float
    fragmented_cjk_ratio: float
    needs_ocr: bool
    warnings: list[str]

    def to_dict(self) -> dict[str, object]:
        return asdict(self)


def load_thresholds(path: Path = QUALITY_CONFIG_PATH) -> tuple[int, dict[str, float]]:
    try:
        payload = json.loads(path.read_text(encoding="utf-8"))
    except (OSError, json.JSONDecodeError) as exc:
        raise RuntimeError(f"cannot read ingestion quality config: {exc}") from exc
    if not isinstance(payload, dict):
        raise RuntimeError("ingestion_quality.json root must be an object")
    thresholds = payload.get("thresholds")
    if payload.get("schema_version") != 1 or not isinstance(payload.get("quality_version"), int) or not isinstance(thresholds, dict):
        raise RuntimeError("ingestion_quality.json has an invalid schema")
    required = (
        "minimum_unit_coverage_pass",
        "minimum_unit_coverage_nonblocked",
        "minimum_average_chars_per_unit",
        "minimum_text_chars_pass",
        "maximum_garbled_ratio",
        "maximum_fragmented_cjk_ratio",
    )
    try:
        normalized = {key: float(thresholds[key]) for key in required}
    except (KeyError, TypeError, ValueError) as exc:
        raise RuntimeError(f"ingestion_quality.json has invalid thresholds: {exc}") from exc
    return int(payload["quality_version"]), normalized


def garbled_ratio(text: str) -> float:
    if not text:
        return 0.0
    suspicious = 0
    for char in text:
        category = unicodedata.category(char)
        if char == "\ufffd" or (category == "Cc" and char not in "\n\r\t"):
            suspicious += 1
    return round(suspicious / len(text), 6)


def normalize_pdf_text(text: str) -> str:
    return re.sub(r"(?<=[\u3400-\u9fff])[ \t]+(?=[\u3400-\u9fff])", "", text)


def fragmented_cjk_ratio(text: str) -> float:
    cjk_count = len(re.findall(r"[\u3400-\u9fff]", text))
    if not cjk_count:
        return 0.0
    fragments = len(re.findall(r"[\u3400-\u9fff][ \t]+(?=[\u3400-\u9fff])", text))
    return round(fragments / cjk_count, 6)


def classify_quality(
    *,
    mode: str,
    total_units: int,
    extracted_texts: list[str],
    quality_version: int,
    thresholds: dict[str, float],
) -> SourceQuality:
    nonempty = [text.strip() for text in extracted_texts if text.strip()]
    extracted_units = len(nonempty)
    total_units = max(total_units, extracted_units)
    empty_units = max(total_units - extracted_units, 0)
    combined = "\n".join(nonempty)
    chars = len(combined)
    coverage = round(extracted_units / total_units, 4) if total_units else 0.0
    average = round(chars / extracted_units, 2) if extracted_units else 0.0
    corruption = garbled_ratio(combined)
    fragmentation = fragmented_cjk_ratio(combined)
    warnings: list[str] = []

    if chars == 0:
        warnings.append("未提取到可用文本。")
    if total_units and coverage < thresholds["minimum_unit_coverage_pass"]:
        warnings.append(f"结构单元提取覆盖率仅为 {coverage:.0%}。")
    if extracted_units and mode in {"pdf", "pptx"} and average < thresholds["minimum_average_chars_per_unit"]:
        warnings.append(f"每个已提取单元平均仅 {average:.0f} 个字符。")
    if mode in {"text", "docx"} and chars < thresholds["minimum_text_chars_pass"]:
        warnings.append(f"提取文本仅 {chars} 个字符。")
    if corruption > thresholds["maximum_garbled_ratio"]:
        warnings.append(f"疑似乱码比例为 {corruption:.2%}。")
    if fragmentation > thresholds["maximum_fragmented_cjk_ratio"]:
        warnings.append(f"中文字符间异常空格比例为 {fragmentation:.2%}。")

    blocked = (
        chars == 0
        or corruption > thresholds["maximum_garbled_ratio"]
        or (total_units > 0 and coverage < thresholds["minimum_unit_coverage_nonblocked"])
    )
    status = "blocked" if blocked else "review" if warnings else "pass"
    needs_ocr = mode == "pdf" and (chars == 0 or coverage < thresholds["minimum_unit_coverage_nonblocked"])
    return SourceQuality(
        quality_version=quality_version,
        extract_mode=mode,
        status=status,
        total_units=total_units,
        extracted_units=extracted_units,
        empty_units=empty_units,
        extracted_chars=chars,
        unit_coverage=coverage,
        average_chars_per_extracted_unit=average,
        garbled_ratio=corruption,
        fragmented_cjk_ratio=fragmentation,
        needs_ocr=needs_ocr,
        warnings=warnings,
    )


def audit_source(path: Path, config_path: Path = QUALITY_CONFIG_PATH) -> SourceQuality:
    source = path.expanduser().resolve()
    if not source.is_file():
        raise FileNotFoundError(f"source file does not exist: {source}")
    quality_version, thresholds = load_thresholds(config_path)
    extension = source.suffix.lower()
    if extension == ".pdf":
        try:
            reader = PdfReader(str(source))
        except Exception:
            return classify_quality(
                mode="pdf", total_units=1, extracted_texts=[], quality_version=quality_version, thresholds=thresholds
            )
        texts: list[str] = []
        for page in reader.pages:
            try:
                texts.append(normalize_pdf_text(page.extract_text() or ""))
            except Exception:
                texts.append("")
        return classify_quality(
            mode="pdf", total_units=len(reader.pages), extracted_texts=texts, quality_version=quality_version, thresholds=thresholds
        )
    if extension == ".pptx":
        try:
            presentation = Presentation(str(source))
        except Exception:
            return classify_quality(
                mode="pptx", total_units=1, extracted_texts=[], quality_version=quality_version, thresholds=thresholds
            )
        texts = []
        for slide in presentation.slides:
            texts.append("\n".join(str(getattr(shape, "text", "") or "") for shape in slide.shapes))
        return classify_quality(
            mode="pptx", total_units=len(presentation.slides), extracted_texts=texts, quality_version=quality_version, thresholds=thresholds
        )
    if extension == ".docx":
        try:
            document = Document(str(source))
        except Exception:
            return classify_quality(
                mode="docx", total_units=1, extracted_texts=[], quality_version=quality_version, thresholds=thresholds
            )
        texts = ["\n".join(paragraph.text for paragraph in document.paragraphs)]
        return classify_quality(
            mode="docx", total_units=1, extracted_texts=texts, quality_version=quality_version, thresholds=thresholds
        )
    if extension in TEXT_EXTENSIONS:
        try:
            text = source.read_text(encoding="utf-8")
        except UnicodeDecodeError:
            text = source.read_text(encoding="utf-8-sig")
        return classify_quality(mode="text", total_units=1, extracted_texts=[text], quality_version=quality_version, thresholds=thresholds)
    return classify_quality(mode="binary", total_units=1, extracted_texts=[], quality_version=quality_version, thresholds=thresholds)


def render_text(quality: SourceQuality) -> str:
    lines = [
        "Source Extraction Quality",
        f"status={quality.status} mode={quality.extract_mode}",
        f"units={quality.extracted_units}/{quality.total_units} coverage={quality.unit_coverage:.0%}",
        f"chars={quality.extracted_chars} average_chars={quality.average_chars_per_extracted_unit:.0f}",
        f"needs_ocr={str(quality.needs_ocr).lower()} garbled_ratio={quality.garbled_ratio:.2%}",
        f"fragmented_cjk_ratio={quality.fragmented_cjk_ratio:.2%}",
    ]
    lines.extend(f"warning: {warning}" for warning in quality.warnings)
    return "\n".join(lines)


def main() -> None:
    parser = argparse.ArgumentParser(description="Audit source extraction quality without printing source content.")
    parser.add_argument("--source", required=True)
    parser.add_argument("--format", choices=["text", "json"], default="text")
    parser.add_argument("--output", default="")
    args = parser.parse_args()
    try:
        quality = audit_source(Path(args.source))
    except (OSError, RuntimeError, ValueError) as exc:
        raise SystemExit(str(exc)) from exc
    content = json.dumps(quality.to_dict(), ensure_ascii=False, indent=2) if args.format == "json" else render_text(quality)
    if args.output:
        output = Path(args.output).expanduser().resolve()
        output.parent.mkdir(parents=True, exist_ok=True)
        output.write_text(content.rstrip() + "\n", encoding="utf-8")
        print(output)
    else:
        print(content)


if __name__ == "__main__":
    main()
