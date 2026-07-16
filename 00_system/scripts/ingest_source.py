from __future__ import annotations

import argparse
import hashlib
import re
import shutil
import subprocess
import sys
from dataclasses import dataclass
from pathlib import Path

from docx import Document
from pptx import Presentation
from pypdf import PdfReader

from create_page import ensure_project
from source_quality import SourceQuality, audit_source, normalize_pdf_text
from wiki_lib import (
    SCRIPT_DIR,
    VAULT_ROOT,
    ai_access_exclusion_reason,
    append_log,
    load_page,
    normalize_tags,
    render_template,
    slugify,
    today_iso,
    update_page_frontmatter,
    write_text,
)

TEMPLATE_DIR = VAULT_ROOT / "00_system" / "templates"
TEXT_EXTENSIONS = {".md", ".markdown", ".txt", ".py", ".js", ".ts", ".json", ".yaml", ".yml", ".html", ".css", ".java", ".go", ".rs", ".sql"}
IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".webp", ".bmp", ".gif", ".tif", ".tiff"}
AUDIO_EXTENSIONS = {".mp3", ".wav", ".m4a", ".aac", ".flac", ".ogg", ".opus"}
VIDEO_EXTENSIONS = {".mp4", ".mov", ".avi", ".mkv", ".webm", ".wmv", ".m4v"}
STRUCTURED_DOCUMENT_MODES = {"text", "docx", "pptx", "pdf"}
SECTION_EXCERPT_LIMIT_CHARS = 1200
HEADING_RE = re.compile(
    r"^(#{1,6}\s+.+|第[一二三四五六七八九十百千万0-9]+[章节篇部分].*|chapter\s+\d+.*|\d+(?:\.\d+){0,4}\.?\s+.+)$",
    re.IGNORECASE,
)
NUMBERED_DOCUMENT_HEADING_RE = re.compile(
    r"^(第[一二三四五六七八九十百千万0-9]+[章节篇部分].*|chapter\s+\d+.*|\d{2}\.?\s+.+)$",
    re.IGNORECASE,
)
APPENDIX_HEADING_RE = re.compile(r"^appendix\s+[a-z0-9].*$", re.IGNORECASE)
LETTERED_CJK_APPENDIX_RE = re.compile(r"^[A-C](?:[.、])?\s+.*[\u3400-\u9fff].*$")
TOC_MARKERS = {"目录", "目次", "table of contents", "contents"}


@dataclass
class SourceBlock:
    ref: str
    label: str
    text: str


@dataclass
class SourceSection:
    title: str
    refs: list[str]
    text: str


def ensure_project_layout(project_name: str) -> Path:
    return ensure_project(project_name, tags=[], status="活跃", summary=f"{project_name} 的项目知识库。")


def ensure_multimodal_scratch_dirs() -> None:
    for rel in (
        Path("01_inbox/scratch/ocr"),
        Path("01_inbox/scratch/transcripts"),
        Path("01_inbox/scratch/keyframes"),
        Path("01_inbox/scratch/summaries"),
        Path("01_inbox/scratch/extracted"),
    ):
        (VAULT_ROOT / rel).mkdir(parents=True, exist_ok=True)


def detect_media_type(path: Path) -> str:
    ext = path.suffix.lower()
    if ext in IMAGE_EXTENSIONS:
        return "image"
    if ext in AUDIO_EXTENSIONS:
        return "audio"
    if ext in VIDEO_EXTENSIONS:
        return "video"
    return "document"


def personal_raw_dir_for_media_type(media_type: str) -> Path:
    base = VAULT_ROOT / "01_inbox" / "raw"
    if media_type == "image":
        return base / "personal" / "images"
    if media_type == "audio":
        return base / "personal" / "audio"
    if media_type == "video":
        return base / "personal" / "video"
    return base


def copy_source(source_file: Path, project_name: str | None, media_type: str) -> Path:
    if project_name:
        destination_dir = ensure_project_layout(project_name) / "sources"
    else:
        destination_dir = personal_raw_dir_for_media_type(media_type)
        destination_dir.mkdir(parents=True, exist_ok=True)

    destination = destination_dir / source_file.name
    if source_file.resolve() != destination.resolve():
        shutil.copy2(source_file, destination)
    return destination


def file_sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def render_source_note(
    title: str,
    domain: str,
    project_slug: str,
    tags: list[str],
    summary: str,
    source_path: str,
    source_hash: str,
    media_type: str,
    extract_mode: str,
    parse_status: str,
    ingest_status: str,
) -> str:
    content = render_template(
        TEMPLATE_DIR / "source-note.md",
        {
            "title": title,
            "type": "来源",
            "domain": domain,
            "project": project_slug,
            "status": "已收录",
            "tags": ", ".join(tags),
            "updated": today_iso(),
            "summary": summary,
        },
    )
    replacements = {
        "source_path:": f"source_path: {source_path}",
        "source_hash:": f"source_hash: {source_hash}",
        "media_type: document": f"media_type: {media_type}",
        "extract_mode:": f"extract_mode: {extract_mode}",
        "parse_status: 已提取": f"parse_status: {parse_status}",
        "ingest_status: 已登记": f"ingest_status: {ingest_status}",
        "review_due:": f"review_due: {today_iso()}",
        "last_parse_attempt:": f"last_parse_attempt: {today_iso()}",
    }
    lines = []
    for line in content.splitlines():
        lines.append(replacements.get(line.strip(), line))
    return "\n".join(lines) + "\n"


def extract_text_from_pptx(path: Path) -> str:
    presentation = Presentation(str(path))
    chunks: list[str] = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        slide_lines: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = (shape.text or "").strip()
                if text:
                    slide_lines.append(text)
        if slide_lines:
            chunks.append(f"## 第 {slide_index} 页\n\n" + "\n\n".join(slide_lines))
    return "\n\n".join(chunks).strip()


def extract_text_from_docx(path: Path) -> str:
    try:
        result = subprocess.run(
            ["pandoc", str(path), "-t", "markdown"],
            check=False,
            capture_output=True,
            text=True,
            encoding="utf-8",
        )
        if result.returncode == 0 and result.stdout.strip():
            return result.stdout.strip()
    except OSError:
        pass

    try:
        document = Document(str(path))
    except Exception:
        return ""
    lines: list[str] = []
    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if not text:
            continue
        style_name = str(getattr(paragraph.style, "name", "") or "")
        if style_name.startswith("Heading "):
            level = style_name.removeprefix("Heading ").strip()
            prefix = "#" * int(level) if level.isdigit() else "##"
            lines.append(f"{prefix} {text}")
        else:
            lines.append(text)
    return "\n\n".join(lines).strip()


def extract_text_from_pdf(path: Path) -> str:
    try:
        reader = PdfReader(str(path))
    except Exception:
        return ""

    chunks: list[str] = []
    for page_index, page in enumerate(reader.pages, start=1):
        try:
            text = normalize_pdf_text(page.extract_text() or "").strip()
        except Exception:
            text = ""
        if text:
            chunks.append(f"## 第 {page_index} 页\n\n{text}")
    return "\n\n".join(chunks).strip()


def extract_text_content(path: Path) -> tuple[str, str]:
    ext = path.suffix.lower()
    if ext in TEXT_EXTENSIONS:
        try:
            return path.read_text(encoding="utf-8"), "text"
        except UnicodeDecodeError:
            try:
                return path.read_text(encoding="utf-8-sig"), "text"
            except UnicodeDecodeError:
                return "", "binary"
    if ext == ".docx":
        return extract_text_from_docx(path), "docx"
    if ext == ".pptx":
        return extract_text_from_pptx(path), "pptx"
    if ext == ".pdf":
        return extract_text_from_pdf(path), "pdf"
    return "", "binary"


def write_extracted_text(title: str, extracted_text: str) -> str:
    if not extracted_text.strip():
        return ""
    output = VAULT_ROOT / "01_inbox" / "scratch" / "extracted" / f"{slugify(title)}.txt"
    write_text(output, extracted_text)
    return output.relative_to(VAULT_ROOT).as_posix()


def block_ref(extract_mode: str, number: str) -> str:
    if extract_mode == "pdf":
        return f"p.{number}"
    if extract_mode == "pptx":
        return f"slide.{number}"
    return f"section.{number}"


def parse_numbered_blocks(extracted_text: str, extract_mode: str) -> list[SourceBlock]:
    blocks: list[SourceBlock] = []
    pattern = re.compile(r"^## 第 (\d+) 页\s*$", re.MULTILINE)
    matches = list(pattern.finditer(extracted_text))
    for index, match in enumerate(matches):
        start = match.end()
        end = matches[index + 1].start() if index + 1 < len(matches) else len(extracted_text)
        number = match.group(1)
        text = extracted_text[start:end].strip()
        if text:
            ref = block_ref(extract_mode, number)
            label = f"第 {number} 页" if extract_mode == "pdf" else f"第 {number} 页/Slide"
            blocks.append(SourceBlock(ref=ref, label=label, text=text))
    return blocks


def parse_heading_blocks(extracted_text: str) -> list[SourceBlock]:
    lines = extracted_text.splitlines()
    blocks: list[SourceBlock] = []
    current_title = "开头"
    current_lines: list[str] = []
    section_index = 1

    for line in lines:
        stripped = line.strip()
        if stripped and HEADING_RE.match(stripped):
            if current_lines:
                blocks.append(
                    SourceBlock(
                        ref=block_ref("text", str(section_index)),
                        label=current_title,
                        text="\n".join(current_lines).strip(),
                    )
                )
                section_index += 1
                current_lines = []
            current_title = stripped
        current_lines.append(line)

    if current_lines:
        blocks.append(
            SourceBlock(
                ref=block_ref("text", str(section_index)),
                label=current_title,
                text="\n".join(current_lines).strip(),
            )
        )
    return [block for block in blocks if block.text.strip()]


def extract_title_from_text(text: str, fallback: str) -> str:
    for raw_line in text.splitlines():
        line = raw_line.strip().lstrip("#").strip()
        if not line:
            continue
        if len(line) <= 80:
            return line
        return line[:80]
    return fallback


def group_blocks(blocks: list[SourceBlock], *, max_chars: int = 9000, max_refs: int = 6) -> list[SourceSection]:
    if not blocks:
        return []

    has_heading_like_blocks = any(HEADING_RE.match(block.label.strip()) for block in blocks)
    sections: list[SourceSection] = []

    if has_heading_like_blocks:
        for block in blocks:
            title = extract_title_from_text(block.label, block.label)
            sections.append(SourceSection(title=title, refs=[block.ref], text=block.text.strip()))
        return sections

    current: list[SourceBlock] = []
    current_size = 0
    for block in blocks:
        should_flush = current and (current_size + len(block.text) > max_chars or len(current) >= max_refs)
        if should_flush:
            title = build_section_title(current)
            sections.append(SourceSection(title=title, refs=[item.ref for item in current], text="\n\n".join(item.text for item in current).strip()))
            current = []
            current_size = 0
        current.append(block)
        current_size += len(block.text)

    if current:
        title = build_section_title(current)
        sections.append(SourceSection(title=title, refs=[item.ref for item in current], text="\n\n".join(item.text for item in current).strip()))
    return sections


def detected_block_heading(block: SourceBlock) -> str:
    lines = [raw_line.strip() for raw_line in block.text.splitlines() if raw_line.strip()]
    for index, raw_line in enumerate(lines):
        line = raw_line.lstrip("#").strip()
        if not line or len(line) > 80:
            continue
        is_numbered = bool(NUMBERED_DOCUMENT_HEADING_RE.match(line))
        is_appendix = bool(APPENDIX_HEADING_RE.match(line) or LETTERED_CJK_APPENDIX_RE.match(line))
        if not is_numbered and not is_appendix:
            continue
        if re.search(r"[。！？.!?]$", line):
            continue
        next_line = lines[index + 1].lstrip("#").strip() if index + 1 < len(lines) else ""
        next_is_translation = (
            len(next_line) <= 80
            and bool(re.search(r"[A-Za-z]{3}", next_line))
            and not bool(re.search(r"[\u3400-\u9fff]", next_line))
            and not NUMBERED_DOCUMENT_HEADING_RE.match(next_line)
            and not APPENDIX_HEADING_RE.match(next_line)
            and not LETTERED_CJK_APPENDIX_RE.match(next_line)
        )
        if index <= 1 or next_is_translation:
            return line
    return ""


def split_front_matter_blocks(blocks: list[SourceBlock]) -> tuple[list[SourceBlock], list[SourceBlock]]:
    toc_index = next(
        (
            index
            for index, block in enumerate(blocks)
            if any(line.strip().lower() in TOC_MARKERS for line in block.text.splitlines())
        ),
        None,
    )
    if toc_index is None:
        return [], blocks

    for index in range(toc_index + 1, len(blocks)):
        block = blocks[index]
        heading = detected_block_heading(block)
        has_body_sentence = bool(re.search(r"[。！？.!?]", block.text))
        if heading and has_body_sentence:
            return blocks[:index], blocks[index:]
    return [], blocks


def label_section_continuations(sections: list[SourceSection]) -> list[SourceSection]:
    previous_title = ""
    continuation = 1
    for section in sections:
        original_title = section.title
        if original_title == previous_title:
            continuation += 1
            section.title = f"{original_title}（续 {continuation}）"
        else:
            previous_title = original_title
            continuation = 1
    return sections


def group_numbered_blocks(blocks: list[SourceBlock], *, max_chars: int = 9000, max_refs: int = 6) -> list[SourceSection]:
    if not blocks:
        return []
    front_matter, body_blocks = split_front_matter_blocks(blocks)
    headings = [detected_block_heading(block) for block in body_blocks]
    if not any(headings):
        return group_blocks(blocks, max_chars=max_chars, max_refs=max_refs)

    sections: list[SourceSection] = []
    if front_matter:
        sections.append(
            SourceSection(
                title="封面与目录",
                refs=[item.ref for item in front_matter],
                text="\n\n".join(item.text for item in front_matter).strip(),
            )
        )
    body_sections: list[SourceSection] = []
    current: list[SourceBlock] = []
    current_heading = ""
    current_size = 0
    for block, heading in zip(body_blocks, headings):
        starts_new_heading = bool(heading and current and heading != current_heading)
        exceeds_limit = bool(current and (current_size + len(block.text) > max_chars or len(current) >= max_refs))
        if starts_new_heading or exceeds_limit:
            body_sections.append(
                SourceSection(
                    title=current_heading or build_section_title(current),
                    refs=[item.ref for item in current],
                    text="\n\n".join(item.text for item in current).strip(),
                )
            )
            current = []
            current_size = 0
        if heading:
            current_heading = heading
        current.append(block)
        current_size += len(block.text)
    if current:
        body_sections.append(
            SourceSection(
                title=current_heading or build_section_title(current),
                refs=[item.ref for item in current],
                text="\n\n".join(item.text for item in current).strip(),
            )
        )
    sections.extend(label_section_continuations(body_sections))
    return sections


def build_section_title(blocks: list[SourceBlock]) -> str:
    if not blocks:
        return "未命名章节"
    first = blocks[0]
    detected = extract_title_from_text(first.text, first.label)
    if len(blocks) == 1:
        return detected
    return f"{detected}（{blocks[0].ref}-{blocks[-1].ref}）"


def build_source_sections(extracted_text: str, extract_mode: str) -> list[SourceSection]:
    if not extracted_text.strip() or extract_mode not in STRUCTURED_DOCUMENT_MODES:
        return []
    if extract_mode in {"pdf", "pptx"}:
        blocks = parse_numbered_blocks(extracted_text, extract_mode)
        return group_numbered_blocks(blocks)
    else:
        blocks = parse_heading_blocks(extracted_text)
    return group_blocks(blocks)


def source_ref_type(extract_mode: str) -> str:
    if extract_mode == "pdf":
        return "page"
    if extract_mode == "pptx":
        return "slide"
    if extract_mode == "docx":
        return "heading"
    return "section"


def page_ref_label(refs: list[str]) -> str:
    if not refs:
        return "-"
    if len(refs) == 1:
        return refs[0]
    return f"{refs[0]} - {refs[-1]}"


def bullet_points_from_text(text: str, limit: int = 5) -> str:
    points: list[str] = []
    for raw_line in text.splitlines():
        line = raw_line.strip().lstrip("-*").strip()
        if not line or line.startswith("```") or len(line) < 8:
            continue
        points.append(f"- {line[:160]}")
        if len(points) >= limit:
            break
    return "\n".join(points) if points else "- 待人工复核。"


def content_lines(text: str) -> list[str]:
    lines: list[str] = []
    for raw_line in text.splitlines():
        line = raw_line.strip().lstrip("-*").strip()
        if not line or line.startswith("```"):
            continue
        lines.append(line)
    return lines


def section_theme_from_text(title: str, text: str) -> str:
    for line in content_lines(text):
        cleaned = line.lstrip("#").strip()
        if 8 <= len(cleaned) <= 140:
            return f"- {cleaned}"
    return f"- {title}"


def key_concepts_from_text(text: str, limit: int = 6) -> str:
    candidates: list[str] = []
    for line in content_lines(text):
        if len(line) > 120:
            continue
        if any(token in line for token in ("定义", "概念", "架构", "模型", "原则", "机制", "流程", "组件", "模块", "接口")):
            candidates.append(f"- {line}")
        if len(candidates) >= limit:
            break
    return "\n".join(candidates) if candidates else "- 待人工复核：本节未自动识别出稳定概念。"


def key_facts_from_text(text: str, limit: int = 6) -> str:
    candidates: list[str] = []
    for line in content_lines(text):
        if len(line) > 180:
            continue
        has_fact_marker = any(char.isdigit() for char in line) or any(
            token in line for token in ("必须", "支持", "包含", "用于", "通过", "需要", "限制", "依赖", "默认")
        )
        if has_fact_marker:
            candidates.append(f"- {line}")
        if len(candidates) >= limit:
            break
    return "\n".join(candidates) if candidates else "- 待人工复核：本节未自动识别出明确事实。"


def process_steps_from_text(text: str, limit: int = 6) -> str:
    candidates: list[str] = []
    step_re = re.compile(r"^(\d+[.)、]|第[一二三四五六七八九十]+步|step\s+\d+)", re.IGNORECASE)
    for line in content_lines(text):
        if len(line) > 180:
            continue
        if step_re.match(line) or any(token in line for token in ("首先", "然后", "接着", "最后", "步骤", "流程", "执行", "运行")):
            candidates.append(f"- {line}")
        if len(candidates) >= limit:
            break
    return "\n".join(candidates) if candidates else "- 本节未自动识别出操作流程；如用于教学或执行，请人工补充。"


def follow_up_questions_for(section_title: str) -> str:
    return "\n".join(
        [
            f"- `{section_title}` 中哪些结论值得提升为正式知识页？",
            "- 本节是否包含需要回到原文核对的参数、步骤或限制？",
            "- 本节内容应该沉淀到个人、项目、共享还是输出层？",
        ]
    )


def recommended_targets_for(domain: str, project_slug: str) -> list[str]:
    if domain == "项目" and project_slug:
        return [f"file-project:{project_slug}", "review-shared"]
    return ["review-personal", "review-shared", "review-output"]


def promotion_candidates_for(domain: str, project_slug: str) -> str:
    targets = recommended_targets_for(domain, project_slug)
    labels = {
        "review-personal": "个人层：长期认知、偏好、学习笔记。",
        "review-shared": "共享层：可跨项目复用的方法、模式、提示词、架构经验。",
        "review-output": "输出层：一次性分析、课程材料、复盘报告。",
    }
    lines: list[str] = []
    for target in targets:
        if target.startswith("file-project:"):
            lines.append(f"- `{target}`：项目事实、项目决策、项目任务或项目来源。")
        else:
            lines.append(f"- `{target}`：{labels.get(target, '待人工复核。')}")
    return "\n".join(lines)


def excerpt_from_text(text: str, limit: int = SECTION_EXCERPT_LIMIT_CHARS) -> str:
    cleaned = "\n".join(line.rstrip() for line in text.strip().splitlines() if line.strip())
    if not cleaned:
        return "- 未提取到可用正文。"
    if len(cleaned) > limit:
        notice = "\n\n> 摘录已截断；完整提取文本见 scratch/extracted。"
        return cleaned[: max(limit - len(notice), 0)] + notice
    return cleaned


def routing_report_for(domain: str, project_slug: str) -> str:
    if domain == "项目":
        return "\n".join(
            [
                f"- 优先沉淀到项目层：`20_projects/active/{project_slug}/`。",
                "- 项目事实更新 `概览.md`、`架构.md`、`决策.md`、`任务.md` 或 `来源.md`。",
                "- 跨项目可复用方法再提升到 `30_shared/`。",
            ]
        )
    return "\n".join(
        [
            "- 个人长期认知进入 `10_personal/`。",
            "- 一次性分析进入 `40_outputs/`。",
            "- 可跨项目复用的方法、提示词、架构经验进入 `30_shared/`。",
        ]
    )


def wiki_link_from_rel(rel_path: str, label: str) -> str:
    return f"[[{Path(rel_path).with_suffix('').as_posix()}|{label}]]"


def render_quoted_list(items: list[str]) -> str:
    return ", ".join(f'"{item}"' for item in items)


def previous_generated_sections(map_path: Path, section_dir: Path) -> list[Path]:
    if not map_path.exists():
        return []
    metadata = load_page(map_path)["frontmatter"]
    derived_sections = metadata.get("derived_sections", [])
    if not isinstance(derived_sections, list):
        return []
    expected_parent = section_dir.resolve()
    paths: list[Path] = []
    for rel_path in derived_sections:
        if not isinstance(rel_path, str):
            continue
        path = (VAULT_ROOT / rel_path).resolve()
        if path.parent == expected_parent:
            paths.append(path)
    return paths


def create_document_derivatives(
    *,
    title: str,
    domain: str,
    project_slug: str,
    tags: list[str],
    source_path: str,
    source_hash: str,
    extract_mode: str,
    extracted_text: str,
    note_path: Path,
    note_dir: Path,
    quality: SourceQuality | None,
) -> list[Path]:
    sections = build_source_sections(extracted_text, extract_mode)
    if not sections:
        return []

    base_slug = slugify(title)
    section_dir = note_dir / f"{base_slug}-sections"
    section_dir.mkdir(parents=True, exist_ok=True)
    map_path = note_dir / f"{base_slug}-document-map.md"
    note_rel = note_path.relative_to(VAULT_ROOT).as_posix()
    map_rel = map_path.relative_to(VAULT_ROOT).as_posix()
    source_note_link = wiki_link_from_rel(note_rel, "来源笔记")
    document_map_link = wiki_link_from_rel(map_rel, "文档地图")
    section_paths = [section_dir / f"{index:02d}-{slugify(section.title)}.md" for index, section in enumerate(sections, start=1)]
    obsolete_section_paths = set(previous_generated_sections(map_path, section_dir)) - {path.resolve() for path in section_paths}

    derived_section_rels = [path.relative_to(VAULT_ROOT).as_posix() for path in section_paths]
    section_index_lines = []
    for index, (section, section_path) in enumerate(zip(sections, section_paths), start=1):
        rel = section_path.relative_to(VAULT_ROOT).as_posix()
        section_index_lines.append(f"- {index:02d}. {wiki_link_from_rel(rel, section.title)} | `{page_ref_label(section.refs)}`")

    quality_lines = [
        f"- 提取方式: `{extract_mode}`",
        f"- 结构段落数: {len(sections)}",
        f"- 引用粒度: `{source_ref_type(extract_mode)}`",
        f"- 完整提取文本: `01_inbox/scratch/extracted/{slugify(title)}.txt`",
        *quality_report_lines(quality),
    ]
    content = render_template(
        TEMPLATE_DIR / "source-document-map.md",
        {
            "title": f"{title} - 文档地图",
            "domain": domain,
            "project": project_slug,
            "tags": ", ".join(tags),
            "updated": today_iso(),
            "summary": f"{title} 的结构地图和章节索引。",
            "source_note": note_rel,
            "source_path": source_path,
            "source_hash": source_hash,
            "extract_mode": extract_mode,
            "source_ref_type": source_ref_type(extract_mode),
            "section_count": str(len(sections)),
            "source_note_link": source_note_link,
            "section_index": "\n".join(section_index_lines),
            "quality_report": "\n".join(quality_lines),
            "routing_report": routing_report_for(domain, project_slug),
        },
    )
    write_text(map_path, content)
    map_updates: dict[str, object] = {"derived_sections": derived_section_rels}
    if quality is not None:
        map_updates.update(quality_frontmatter(quality))
    update_page_frontmatter(map_path, map_updates)

    for index, (section, section_path) in enumerate(zip(sections, section_paths), start=1):
        section_title = f"{title} - {index:02d} {section.title}"
        recommended_targets = recommended_targets_for(domain, project_slug)
        content = render_template(
            TEMPLATE_DIR / "source-section-note.md",
            {
                "title": section_title,
                "domain": domain,
                "project": project_slug,
                "tags": ", ".join(tags),
                "updated": today_iso(),
                "summary": f"{title} 的章节笔记：{section.title}。",
                "source_note": note_rel,
                "source_path": source_path,
                "source_hash": source_hash,
                "source_refs": render_quoted_list(section.refs),
                "document_map": map_rel,
                "excerpt_limit_chars": str(SECTION_EXCERPT_LIMIT_CHARS),
                "recommended_targets": render_quoted_list(recommended_targets),
                "source_note_link": source_note_link,
                "document_map_link": document_map_link,
                "source_ref_label": page_ref_label(section.refs),
                "key_points": bullet_points_from_text(section.text),
                "section_theme": section_theme_from_text(section.title, section.text),
                "key_concepts": key_concepts_from_text(section.text),
                "key_facts": key_facts_from_text(section.text),
                "process_steps": process_steps_from_text(section.text),
                "structured_excerpt": excerpt_from_text(section.text),
                "follow_up_questions": follow_up_questions_for(section.title),
                "promotion_candidates": promotion_candidates_for(domain, project_slug),
                "routing_report": routing_report_for(domain, project_slug),
            },
        )
        write_text(section_path, content)
    for obsolete_path in obsolete_section_paths:
        obsolete_path.unlink(missing_ok=True)
    return [map_path, *section_paths]


def detect_parse_status(
    media_type: str, extracted_text: str, extract_mode: str, quality: SourceQuality | None = None
) -> str:
    if media_type == "document":
        if quality is not None and quality.status == "blocked":
            return "需OCR" if quality.needs_ocr else "提取失败"
        if quality is not None and quality.status == "review":
            return "待复核"
        return "已提取" if extract_mode != "binary" else "待处理"
    return "待处理"


def bool_literal(value: bool) -> str:
    return "true" if value else "false"


def quality_frontmatter(quality: SourceQuality) -> dict[str, object]:
    return {
        "ingest_quality_version": quality.quality_version,
        "quality_status": quality.status,
        "quality_total_units": quality.total_units,
        "quality_extracted_units": quality.extracted_units,
        "quality_empty_units": quality.empty_units,
        "quality_extracted_chars": quality.extracted_chars,
        "quality_unit_coverage": quality.unit_coverage,
        "quality_average_chars_per_unit": quality.average_chars_per_extracted_unit,
        "quality_garbled_ratio": quality.garbled_ratio,
        "quality_fragmented_cjk_ratio": quality.fragmented_cjk_ratio,
        "needs_ocr": quality.needs_ocr,
        "quality_warnings": quality.warnings,
    }


def quality_report_lines(quality: SourceQuality | None) -> list[str]:
    if quality is None:
        return ["- 抽取质量门禁: 不适用。"]
    lines = [
        f"- 抽取质量门禁: `{quality.status}`",
        f"- 结构单元覆盖: {quality.extracted_units}/{quality.total_units} ({quality.unit_coverage:.0%})",
        f"- 提取字符数: {quality.extracted_chars}",
        f"- 已提取单元平均字符数: {quality.average_chars_per_extracted_unit:.0f}",
        f"- 疑似乱码比例: {quality.garbled_ratio:.2%}",
        f"- 中文字符异常空格比例: {quality.fragmented_cjk_ratio:.2%}",
        f"- 需要 OCR: `{bool_literal(quality.needs_ocr)}`",
    ]
    lines.extend(f"- 质量警告: {warning}" for warning in quality.warnings)
    return lines


def enrich_source_note(
    content: str,
    extracted_text: str,
    extract_mode: str,
    media_type: str,
    quality: SourceQuality | None = None,
) -> str:
    lines = content.rstrip().splitlines()
    if media_type == "image":
        lines.extend(["", "## 媒体处理", "", "- 媒体类型: 图片", "- 当前阶段: 已入库，待 OCR / caption。"])
    elif media_type == "audio":
        lines.extend(["", "## 媒体处理", "", "- 媒体类型: 语音", "- 当前阶段: 已入库，待转写 / 摘要。"])
    elif media_type == "video":
        lines.extend(["", "## 媒体处理", "", "- 媒体类型: 视频", "- 当前阶段: 已入库，待音轨转写 / 关键帧 / 时间轴摘要。"])
    lines.extend(["", "## 文件识别", "", f"- 识别方式: {extract_mode}"])
    if quality is not None:
        lines.extend(quality_report_lines(quality))
    if extracted_text.strip():
        excerpt = extracted_text.strip()[:1200]
        if len(extracted_text.strip()) > 1200:
            excerpt += "\n\n[摘录已截断；请阅读文档地图和章节笔记。]"
        lines.extend(["", "## 提取文本摘录", "", "```text", excerpt, "```"])
    else:
        lines.extend(["", "## 提取文本摘录", "", "- 当前未提取到正文文本。文件已入库，可后续人工补充或扩展解析能力。"])
    return "\n".join(lines) + "\n"


def append_to_project_sources(project_name: str, note_path: Path, stored_source: Path) -> None:
    source_registry = ensure_project_layout(project_name) / "来源.md"
    if not source_registry.exists():
        return

    relative_source = stored_source.relative_to(VAULT_ROOT).as_posix()
    note_target = note_path.relative_to(source_registry.parent).with_suffix("").as_posix()
    entry = f"- [[{note_target}|{note_path.stem}]] | 原文件: `{relative_source}`"
    text = source_registry.read_text(encoding="utf-8")
    if entry not in text:
        source_registry.write_text(text.rstrip() + "\n" + entry + "\n", encoding="utf-8")


def detect_ingest_status(extracted_text: str, extract_mode: str, quality: SourceQuality | None = None) -> str:
    if quality is not None and quality.status == "blocked":
        return "待人工处理"
    if quality is not None and quality.status == "review":
        return "待复核"
    if extract_mode == "binary":
        return "已登记"
    if extracted_text.strip():
        return "已解析"
    return "已登记"


def main() -> None:
    parser = argparse.ArgumentParser(description="把原始资料复制到知识库并生成来源笔记。")
    parser.add_argument("--source", required=True, help="原始文件路径")
    parser.add_argument("--title", default="", help="来源标题，默认取文件名")
    parser.add_argument("--project", default="", help="所属项目名")
    parser.add_argument("--tags", default="", help="英文逗号分隔标签")
    parser.add_argument("--summary", default="", help="一句话摘要")
    args = parser.parse_args()

    source_file = Path(args.source).expanduser().resolve()
    if not source_file.exists():
        raise SystemExit(f"源文件不存在: {source_file}")

    title = args.title.strip() or source_file.stem
    project_name = args.project.strip() or None
    project_slug = slugify(project_name) if project_name else ""
    tags = normalize_tags(args.tags)
    summary = args.summary.strip() or f"{title} 的来源笔记。"
    media_type = detect_media_type(source_file)
    prospective_source = (
        VAULT_ROOT / "20_projects" / "active" / project_slug / "sources" / source_file.name
        if project_name
        else personal_raw_dir_for_media_type(media_type) / source_file.name
    )
    exclusion = ai_access_exclusion_reason(source_file) or ai_access_exclusion_reason(prospective_source)
    if exclusion:
        raise SystemExit(f"source is blocked by AI access policy: {exclusion}")
    ensure_multimodal_scratch_dirs()
    stored_source = copy_source(source_file, project_name, media_type)
    extracted_text, extract_mode = extract_text_content(stored_source)
    quality = audit_source(stored_source) if media_type == "document" else None
    extracted_text_path = write_extracted_text(title, extracted_text)
    source_hash = file_sha256(stored_source)
    parse_status = detect_parse_status(media_type, extracted_text, extract_mode, quality)
    ingest_status = detect_ingest_status(extracted_text, extract_mode, quality)

    if project_name:
        note_dir = ensure_project_layout(project_name) / "source-notes"
        domain = "项目"
    else:
        note_dir = VAULT_ROOT / "01_inbox" / "clips"
        note_dir.mkdir(parents=True, exist_ok=True)
        domain = "个人"

    note_path = note_dir / f"{slugify(title)}.md"
    content = render_source_note(
        title=title,
        domain=domain,
        project_slug=project_slug,
        tags=tags,
        summary=summary,
        source_path=stored_source.relative_to(VAULT_ROOT).as_posix(),
        source_hash=source_hash,
        media_type=media_type,
        extract_mode=extract_mode,
        parse_status=parse_status,
        ingest_status=ingest_status,
    )
    replacements = {
        "has_ocr_text:": f"has_ocr_text: {bool_literal(bool(extracted_text.strip()) and media_type == 'image')}",
        "has_transcript:": f"has_transcript: {bool_literal(bool(extracted_text.strip()) and media_type in {'audio', 'video'})}",
        "has_keyframes:": f"has_keyframes: {bool_literal(False)}",
    }
    lines = []
    for line in content.splitlines():
        lines.append(replacements.get(line.strip(), line))
    content = "\n".join(lines) + "\n"
    content = enrich_source_note(content, extracted_text, extract_mode, media_type, quality)
    write_text(note_path, content)
    if quality is not None:
        update_page_frontmatter(note_path, quality_frontmatter(quality))
    derived_paths = []
    if quality is None or quality.status != "blocked":
        derived_paths = create_document_derivatives(
            title=title,
            domain=domain,
            project_slug=project_slug,
            tags=tags,
            source_path=stored_source.relative_to(VAULT_ROOT).as_posix(),
            source_hash=source_hash,
            extract_mode=extract_mode,
            extracted_text=extracted_text,
            note_path=note_path,
            note_dir=note_dir,
            quality=quality,
        )
    if derived_paths:
        update_page_frontmatter(
            note_path,
            {
                "derived_pages": [path.relative_to(VAULT_ROOT).as_posix() for path in derived_paths],
                "ingest_status": "待复核" if quality is not None and quality.status == "review" else "已解析",
                "recommended_targets": [f"file-project:{project_slug}"] if project_slug else ["review-personal", "review-shared"],
            },
        )
    if extracted_text_path:
        text = note_path.read_text(encoding="utf-8")
        marker = "## 文件识别\n\n"
        if marker in text and "完整提取文本:" not in text:
            text = text.replace(marker, marker + f"- 完整提取文本: `{extracted_text_path}`\n", 1)
            write_text(note_path, text)

    if project_name:
        append_to_project_sources(project_name, note_path, stored_source)

    details = f"来源文件: {stored_source.relative_to(VAULT_ROOT).as_posix()} | 来源笔记: {note_path.relative_to(VAULT_ROOT).as_posix()}"
    append_log("摄入", title, details)
    subprocess.run([sys.executable, str(SCRIPT_DIR / "rebuild_indexes.py")], check=True)
    print(note_path)


if __name__ == "__main__":
    main()

